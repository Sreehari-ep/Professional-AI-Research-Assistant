from pathlib import Path
import fitz,pandas as pd
from docx import Document
from pptx import Presentation

def extract_text(filepath):
    path=Path(filepath); ext=path.suffix.lower()
    if ext==".pdf": text="\n".join(page.get_text() for page in fitz.open(filepath))
    elif ext==".docx": text="\n".join(p.text for p in Document(filepath).paragraphs if p.text.strip())
    elif ext==".txt": text=path.read_text(encoding="utf-8",errors="ignore")
    elif ext==".csv": text=pd.read_csv(filepath).astype(str).to_string(index=False)
    elif ext in {".xlsx",".xls"}: text=pd.read_excel(filepath).astype(str).to_string(index=False)
    elif ext==".pptx":
        prs=Presentation(filepath)
        text="\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape,"text"))
    else: raise ValueError(f"Unsupported file type: {ext}")
    if not str(text).strip(): raise ValueError("No readable text extracted.")
    return str(text).strip()
